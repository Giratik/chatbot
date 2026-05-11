#backend/file_type_action.py

import filetype
from llm_vision import analyse_image
from paddle_ocr_processor import process_file_with_ocr

import io
import zipfile
import docx
from pptx import Presentation

MIMES_OFFICE = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
}


# --- 1. FONCTIONS SPÉCIALISÉES ---

def traiter_image(fichier_joint, model_vlm):
    image_bytes = fichier_joint.getvalue()
    return analyse_image(image_bytes, prompt="Voici l'image: Rédige une description détaillée de cette image en français.", model=model_vlm)

def traiter_pdf(fichier_bytes, model_vlm):
    print("Analyse du PDF avec OCR...")
    return process_file_with_ocr(fichier_bytes, file_type="application/pdf")

def traiter_texte(fichier_objet, model_vlm):
    fichier_bytes = fichier_objet.read()
    if hasattr(fichier_objet, 'seek'):
        fichier_objet.seek(0)

    print("Analyse du contenu texte...")
    contenu_texte = fichier_bytes.decode('utf-8', errors='replace')

    if contenu_texte.count('\ufffd') > (len(contenu_texte) * 0.1):
        raise ValueError("Probablement un fichier binaire")

    return contenu_texte


# --- HELPERS COMMUNS ---

FORMATS_IMAGE = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'}

def _extraire_images_zip(fichier_bytes: bytes, dossier_interne: str) -> list[bytes]:
    """
    Extrait les images embarquées dans un fichier Office (docx/pptx),
    qui sont en réalité des archives ZIP.
    dossier_interne : 'word/media' pour docx, 'ppt/media' pour pptx.
    """
    images = []
    try:
        with zipfile.ZipFile(io.BytesIO(fichier_bytes)) as z:
            for nom in z.namelist():
                ext = '.' + nom.rsplit('.', 1)[-1].lower() if '.' in nom else ''
                if nom.startswith(dossier_interne) and ext in FORMATS_IMAGE:
                    images.append(z.read(nom))
    except Exception as e:
        print(f"Erreur extraction images ZIP: {e}")
    return images

def _ocr_images(images: list[bytes]) -> str:
    """Passe chaque image par PaddleOCR et concatène les résultats."""
    resultats = []
    for i, img_bytes in enumerate(images, start=1):
        print(f"  OCR image embarquée {i}/{len(images)}...")
        try:
            texte_ocr = process_file_with_ocr(img_bytes, file_type="image/png")
            if texte_ocr and texte_ocr.strip():
                resultats.append(f"[Image {i}]\n{texte_ocr.strip()}")
        except Exception as e:
            print(f"  Erreur OCR image {i}: {e}")
    return '\n\n'.join(resultats)


# --- DOCX ---

def traiter_docx(fichier_objet, model_vlm):
    """Extrait le texte natif + OCR des images embarquées d'un fichier Word."""
    print("Extraction du contenu Word...")
    fichier_bytes = fichier_objet.getvalue()
    parties = []

    # 1. Texte natif (paragraphes + tableaux)
    try:
        doc = docx.Document(io.BytesIO(fichier_bytes))

        paragraphes = [para.text for para in doc.paragraphs if para.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                ligne = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if ligne:
                    paragraphes.append(ligne)

        if paragraphes:
            parties.append("=== Texte du document ===\n" + '\n'.join(paragraphes))

    except Exception as e:
        print(f"Erreur extraction texte DOCX: {e}")

    # 2. OCR sur les images embarquées (dans word/media/)
    images = _extraire_images_zip(fichier_bytes, dossier_interne='word/media')
    if images:
        print(f"{len(images)} image(s) embarquée(s) trouvée(s) dans le Word.")
        texte_ocr = _ocr_images(images)
        if texte_ocr:
            parties.append("=== Texte extrait des images (OCR) ===\n" + texte_ocr)
    else:
        print("Aucune image embarquée dans le Word.")

    return '\n\n'.join(parties) if parties else "⚠️ Impossible de lire le fichier Word."


# --- PPTX ---

def traiter_pptx(fichier_objet, model_vlm):
    """Extrait le texte natif + OCR des images embarquées d'un fichier PowerPoint."""
    print("Extraction du contenu PowerPoint...")
    fichier_bytes = fichier_objet.getvalue()
    parties = []

    # 1. Texte natif (shapes textuelles de chaque slide)
    try:
        prs = Presentation(io.BytesIO(fichier_bytes))
        slides_texte = []

        for i, slide in enumerate(prs.slides, start=1):
            blocs = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if blocs:
                slides_texte.append(f"--- Slide {i} ---\n" + '\n'.join(blocs))

        if slides_texte:
            parties.append("=== Texte des slides ===\n" + '\n\n'.join(slides_texte))

    except Exception as e:
        print(f"Erreur extraction texte PPTX: {e}")

    # 2. OCR sur les images embarquées (dans ppt/media/)
    images = _extraire_images_zip(fichier_bytes, dossier_interne='ppt/media')
    if images:
        print(f"{len(images)} image(s) embarquée(s) trouvée(s) dans le PowerPoint.")
        texte_ocr = _ocr_images(images)
        if texte_ocr:
            parties.append("=== Texte extrait des images (OCR) ===\n" + texte_ocr)
    else:
        print("Aucune image embarquée dans le PowerPoint.")

    return '\n\n'.join(parties) if parties else "⚠️ Impossible de lire le fichier PowerPoint."


def format_non_supporte(fichier_bytes, model_vlm):
    return "⚠️ Ce type de contenu n'est pas encore pris en charge."


# --- 2. DICTIONNAIRE DE ROUTAGE ---

ROUTEUR = {
    "image": traiter_image,
    "pdf":   traiter_pdf,
    "texte": traiter_texte,
    "docx":  traiter_docx,
    "pptx":  traiter_pptx,
}

EXTENSIONS_TEXTE = {"txt", "csv", "md", "json", "xml", "html", "py", "js", "ts"}

# --- 3. CHEF D'ORCHESTRE ---

def analyser_contenu_fichier(uploaded_file, model_vlm):
    fichier_bytes = uploaded_file.getvalue()
    kind = filetype.guess(fichier_bytes)
    categorie = None

    nom = getattr(uploaded_file, 'name', '')
    extension = nom.rsplit('.', 1)[-1].lower() if '.' in nom else ''

    # 1. Extension en priorité pour les formats Office
    if extension in ('docx', 'pptx'):
        categorie = extension

    elif kind is None:
        if extension in EXTENSIONS_TEXTE or not extension:
            try:
                fichier_bytes.decode('utf-8')
                categorie = 'texte'
            except UnicodeDecodeError:
                categorie = 'inconnu'

    else:
        mime = kind.mime
        mime_principal = mime.split('/')[0]
        extension_reelle = kind.extension

        if mime in MIMES_OFFICE:
            # filetype reconnaît le MIME Office complet → on mappe directement
            categorie = MIMES_OFFICE[mime]
        elif mime_principal == 'image':
            categorie = 'image'
        elif extension_reelle == 'pdf':
            categorie = 'pdf'
        elif mime_principal == 'audio':
            categorie = 'audio'

    print(f">>> Catégorie détectée : {categorie}")

    fonction_a_executer = ROUTEUR.get(categorie, format_non_supporte)
    return fonction_a_executer(uploaded_file, model_vlm)